import io
import tempfile
import os

import pandas as pd
import pytest
from pptx import Presentation
from pptx.util import Inches

from tools.t6_pptx_builder import PPTXBuilder


def _make_builder():
    with tempfile.TemporaryDirectory() as d:
        builder = PPTXBuilder(templates_dir=d)
        return builder


def _ai_output():
    return {
        "summary": "Test summary",
        "kpi_evaluation": [],
        "key_achievements": ["Ach 1"],
        "insights": ["Insight 1"],
    }


def test_add_group_slide_original_interface():
    builder = _make_builder()
    builder.start_or_load_report(os.path.join(tempfile.mkdtemp(), "out.pptx"))
    builder.add_title_slide("20260417", 1)
    df = pd.DataFrame({"Week": ["W1", "W2"], "SLA": [0.95, 0.97]})
    builder.add_group_slide(
        sheet_name="MIM",
        ai_output=_ai_output(),
        kpi_rows=[],
        chart_bytes=None,
        df=df,
    )
    assert len(builder.prs.slides) == 2


def test_add_group_slide_with_charts_list():
    builder = _make_builder()
    builder.start_or_load_report(os.path.join(tempfile.mkdtemp(), "out.pptx"))
    builder.add_title_slide("20260417", 1)
    df = pd.DataFrame({"A": [1, 2], "B": [3, 4]})
    builder.add_group_slide(
        sheet_name="MIM",
        ai_output=_ai_output(),
        kpi_rows=[],
        chart_bytes=None,
        charts=[{"chart_id": "chart_0", "png_bytes": None, "position": "right", "title": "T"}],
        summary_mode="ai_write",
        df=df,
    )
    # Slide should exist; data table should have been placed
    slide = builder.prs.slides[-1]
    table_shapes = [s for s in slide.shapes if s.has_table]
    assert len(table_shapes) >= 1


def test_add_data_table_row_limit():
    builder = _make_builder()
    builder.start_or_load_report(os.path.join(tempfile.mkdtemp(), "out.pptx"))
    builder.add_title_slide("20260417", 1)
    slide = builder.prs.slides.add_slide(
        builder.prs.slide_layouts[builder._blank_slide_layout()]
    )
    df = pd.DataFrame({"A": list(range(10)), "B": list(range(10))})
    builder._add_data_table(slide, df, (0.5, 1.5, 5.0, 3.0))
    table_shapes = [s for s in slide.shapes if s.has_table]
    assert len(table_shapes) == 1
    tbl = table_shapes[0].table
    # max 4 data rows + 1 header = 5
    assert tbl.rows.__len__() <= 5


def test_summary_mode_use_excel_label():
    builder = _make_builder()
    builder.start_or_load_report(os.path.join(tempfile.mkdtemp(), "out.pptx"))
    builder.add_title_slide("20260417", 1)
    df = pd.DataFrame({"Week": ["W1", "W2"], "SLA": [0.95, 0.97]})
    builder.add_group_slide(
        sheet_name="NOC",
        ai_output=_ai_output(),
        kpi_rows=[],
        chart_bytes=None,
        summary_mode="use_excel",
        df=df,
    )
    slide = builder.prs.slides[-1]
    all_text = " ".join(
        shape.text_frame.text for shape in slide.shapes
        if hasattr(shape, "text_frame")
    )
    assert "(Team-provided)" in all_text
