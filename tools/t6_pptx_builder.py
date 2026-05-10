import io
import json
import math
import os
import re
from datetime import datetime, date as date_type
from typing import Any, Dict, List, Optional

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE


class PPTXBuilder:
    def __init__(self, templates_dir: Optional[str] = None):
        base_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
        self.templates_dir = templates_dir or os.path.join(base_dir, "templates")
        self.template_path = os.path.join(self.templates_dir, "slide_template.pptx")

        if os.path.exists(self.template_path):
            self.prs = Presentation(self.template_path)
        else:
            os.makedirs(self.templates_dir, exist_ok=True)
            self.prs = Presentation()
            self.prs.save(self.template_path)

        self.prs.slide_width  = Inches(13.33)
        self.prs.slide_height = Inches(7.5)

        self.MUFG_RED   = RGBColor(0xCC, 0x00, 0x00)
        self.MUFG_GRAY  = RGBColor(0x59, 0x57, 0x57)
        self.MUFG_DARK  = RGBColor(0x3A, 0x3A, 0x3A)
        self.MUFG_LGRAY = RGBColor(0xF5, 0xF5, 0xF5)
        self.MUFG_RED2  = RGBColor(0xE6, 0x33, 0x33)
        self.DARK_BLUE  = self.MUFG_RED
        self.MID_BLUE   = self.MUFG_GRAY
        self.TEAL       = RGBColor(0x0F, 0x6E, 0x56)
        self.RED        = self.MUFG_RED
        self.GREEN      = RGBColor(0x1A, 0x73, 0x48)
        self.LIGHT_BG   = self.MUFG_LGRAY

        self.report_date: str = datetime.now().strftime("%Y%m%d")

    def start_or_load_report(self, output_path: str):
        if os.path.exists(output_path):
            self.prs = Presentation(output_path)
            self.prs.slide_width  = Inches(13.33)
            self.prs.slide_height = Inches(7.5)
            return
        if os.path.exists(self.template_path):
            self.prs = Presentation(self.template_path)
        else:
            self.prs = Presentation()
        self.prs.slide_width  = Inches(13.33)
        self.prs.slide_height = Inches(7.5)

    @staticmethod
    def _fmt_date(raw: str) -> str:
        try:
            dt = datetime.strptime(raw, "%Y%m%d")
            return f"{dt.day} {dt.strftime('%b')} {dt.year}"
        except Exception:
            return raw

    def _blank_slide_layout(self):
        for i, layout in enumerate(self.prs.slide_layouts):
            if getattr(layout, "name", "") and "Blank" in layout.name:
                return i
        return len(self.prs.slide_layouts) - 1

    def _clear_slide(self, slide):
        for shape in list(slide.shapes):
            shape.element.getparent().remove(shape.element)

    def _find_group_slide(self, sheet_name: str):
        prefix = f"{sheet_name}  |  "
        for slide in self.prs.slides:
            for shape in slide.shapes:
                if not getattr(shape, "has_text_frame", False):
                    continue
                if (shape.text_frame.text or "").strip().startswith(prefix):
                    return slide
        return None

    # ── Data sanitization ─────────────────────────────────────────────────────

    @staticmethod
    def _sanitize_text_list(items: Any) -> List[str]:
        """Extract clean strings from AI output — handles dicts and raw JSON strings."""
        if not items:
            return []
        _KEYS = ("description", "recommendation", "text", "content",
                 "action", "achievement", "insight", "finding", "summary", "title")
        result = []
        for item in items:
            if isinstance(item, dict):
                text = next((str(item[k]) for k in _KEYS if item.get(k)), None)
                if text is None:
                    text = next((str(v) for v in item.values() if isinstance(v, str)),
                                str(item))
                result.append(text.strip())
            elif isinstance(item, str):
                s = item.strip()
                if s.startswith("{") and s.endswith("}"):
                    try:
                        obj = json.loads(s)
                        text = next((str(obj[k]) for k in _KEYS if obj.get(k)), None)
                        if text is None:
                            text = next(
                                (str(v) for v in obj.values() if isinstance(v, str)), s)
                        result.append(text.strip())
                        continue
                    except Exception:
                        pass
                if s:
                    result.append(s)
            elif item is not None:
                s = str(item).strip()
                if s:
                    result.append(s)
        return result

    # ── Rich-text bullet helper ───────────────────────────────────────────────

    def _apply_rich_bullet_to_para(
        self,
        p,
        text: str,
        bullet_char: str = "▸",
        base_pt: float = 8.0,
        text_color: Optional[RGBColor] = None,
        num_color: Optional[RGBColor] = None,
    ) -> None:
        """Fill paragraph with runs where numbers/percentages are bolded."""
        text_color = text_color or self.MUFG_DARK
        num_color  = num_color  or self.MUFG_RED
        p.space_before = Pt(2)
        p.space_after  = Pt(1)

        full_text = f"{bullet_char}  {text}" if bullet_char else text
        _PAT = re.compile(r"(\b\d[\d,]*(?:\.\d+)?%?|\d+/\d+\b)")
        parts = _PAT.split(full_text)

        for i, part in enumerate(parts):
            if not part:
                continue
            run = p.add_run()
            run.text = part
            run.font.size = Pt(base_pt)
            if _PAT.fullmatch(part):
                run.font.bold = True
                run.font.color.rgb = num_color
            else:
                run.font.bold = bool(bullet_char) and (i == 0)
                run.font.color.rgb = text_color

    # ── Section-height estimator ──────────────────────────────────────────────

    @staticmethod
    def _estimate_section_h(
        body_text: str,
        items: List[str],
        label_h: float = 0.30,
        line_h: float = 0.165,
        chars_per_line: int = 85,
        pad: float = 0.15,
    ) -> float:
        """Estimate shape height in inches from content volume."""
        lines = 0
        if body_text:
            lines += max(1, math.ceil(len(body_text) / chars_per_line))
        for item in items:
            lines += max(1, math.ceil(len(str(item)) / chars_per_line))
        return label_h + lines * line_h + pad

    # ── Display-name mapping: Excel sheet name → PPTX slide header ─────────────
    # Add entries here to decouple the internal sheet name from the visual title.
    _DISPLAY_NAME_MAP: Dict[str, str] = {
        "change management": "Asset and Configuration Management",
        # "Asset and Configuration" maps to itself — no rename needed,
        # but listed here explicitly so intent is clear.
        "asset and configuration": "Asset and Configuration",
    }

    # ── Layout constants ──────────────────────────────────────────────────────
    # 60 / 40 split: left text column is ~60 %, right chart column is ~40 %
    _HDR_H        = 0.62   # header bar height
    _PERF_Y       = 0.62   # performance table top
    _PERF_H       = 0.80   # performance table height
    _CONTENT_Y    = 1.42   # content area starts below perf table
    _KPI_DEF_Y    = 6.50   # KPI definition table top
    _KPI_DEF_H    = 1.00   # KPI definition table height
    _SLIDE_BOTTOM = 6.50   # usable bottom edge
    _LEFT_W       = 8.00   # left (text) column width  — ~60 % of 13.33"
    _RIGHT_X      = 8.15   # right (chart) column left edge
    _CHART_W      = 4.95   # right column width — leaves 0.23" right margin to avoid label clipping
    _CHART_FOOTER = 6.50   # charts end here

    # ── Performance & KPI tables (full-width, unchanged) ─────────────────────

    def _add_performance_table(self, slide, df: Optional[pd.DataFrame],
                                kpi_rows: List[Dict]):
        """Returns the table shape so callers can anchor downstream elements to its bottom."""
        if df is None or df.empty:
            return None
        SKIP_KW = {"comment", "achievement", "weekly", "month", "quarter", "date"}
        status_map: Dict[str, str] = {
            str(r.get("kpi_name", "")): str(r.get("status", "")) for r in kpi_rows}
        last = df.iloc[-1]
        cols_to_show = [str(c) for c in df.columns
                        if not any(kw in str(c).lower() for kw in SKIP_KW)]
        if not cols_to_show:
            return None
        cols_to_show = cols_to_show[:16]
        n_cols = len(cols_to_show)
        tbl_shape = slide.shapes.add_table(
            2, n_cols, Inches(0), Inches(self._PERF_Y),
            self.prs.slide_width, Inches(self._PERF_H))
        tbl = tbl_shape.table
        HDR_ROW_H = 0.35
        tbl.rows[0].height = int(Inches(HDR_ROW_H))
        tbl.rows[1].height = int(Inches(self._PERF_H - HDR_ROW_H))
        col_w = Inches(13.33 / n_cols)
        for ci in range(n_cols):
            tbl.columns[ci].width = col_w

        PASS_GREEN = RGBColor(0x1A, 0x73, 0x48)
        AMBER_COL  = RGBColor(0xFF, 0xC1, 0x07)
        NEUTRAL_BG = self.MUFG_LGRAY

        CELL_MARGIN = int(Inches(0.05))   # minimal padding so larger text doesn't wrap

        for ci, col_str in enumerate(cols_to_show):
            display_name = col_str.split("(")[0].strip() if "(" in col_str else col_str

            # ── Header cell ───────────────────────────────────────────────────
            hdr = tbl.cell(0, ci)
            hdr.text = display_name
            hdr.fill.solid()
            hdr.fill.fore_color.rgb = self.MUFG_RED
            hdr.margin_top    = CELL_MARGIN
            hdr.margin_bottom = CELL_MARGIN
            hdr.margin_left   = CELL_MARGIN
            hdr.margin_right  = CELL_MARGIN
            for p in hdr.text_frame.paragraphs:
                p.font.bold = True
                p.font.size = Pt(9)   # keep small so wrapped headers don't expand the row
                p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                p.alignment = PP_ALIGN.CENTER
            hdr.text_frame.word_wrap = True

            # ── Data cell ─────────────────────────────────────────────────────
            raw_val = last.get(col_str) if hasattr(last, "get") else last[col_str]
            val_cell = tbl.cell(1, ci)
            val_str = ""
            if raw_val is not None and not (isinstance(raw_val, float) and raw_val != raw_val):
                try:
                    fv = float(str(raw_val).replace("%", "").strip())
                    if 0 < fv <= 1.0:
                        val_str = f"{fv*100:.0f}%"
                    elif fv > 1.0 and fv <= 100 and "%" in col_str:
                        val_str = f"{fv:.0f}%"
                    else:
                        val_str = f"{int(fv)}" if fv == int(fv) else f"{fv:.1f}"
                except Exception:
                    val_str = str(raw_val)
            val_cell.text = val_str
            val_cell.margin_top    = CELL_MARGIN
            val_cell.margin_bottom = CELL_MARGIN
            val_cell.margin_left   = CELL_MARGIN
            val_cell.margin_right  = CELL_MARGIN

            status = status_map.get(col_str.split("(")[0].strip(), "")
            if status == "PASS":
                bg, fg = PASS_GREEN, RGBColor(0xFF, 0xFF, 0xFF)
            elif status == "AMBER":
                bg, fg = AMBER_COL, RGBColor(0x33, 0x33, 0x33)
            elif status == "FAIL":
                bg, fg = self.MUFG_RED, RGBColor(0xFF, 0xFF, 0xFF)
            else:
                bg, fg = NEUTRAL_BG, RGBColor(0x33, 0x33, 0x33)
            val_cell.fill.solid()
            val_cell.fill.fore_color.rgb = bg
            for p in val_cell.text_frame.paragraphs:
                p.font.size = Pt(13)
                p.font.bold = True
                p.font.color.rgb = fg
                p.alignment = PP_ALIGN.CENTER

        return tbl_shape

    def _add_kpi_definition_table(self, slide, kpi_definitions: List[Dict]) -> None:
        defs = [d for d in kpi_definitions
                if d.get("kpi_name") and (d.get("definition") or d.get("sub_definition"))]
        if not defs:
            return
        defs = defs[:4]
        n_rows = len(defs) + 1
        tbl_shape = slide.shapes.add_table(
            n_rows, 3, Inches(0), Inches(self._KPI_DEF_Y),
            self.prs.slide_width, Inches(self._KPI_DEF_H))
        tbl = tbl_shape.table
        for ci, w in enumerate([Inches(2.8), Inches(5.4), Inches(5.13)]):
            tbl.columns[ci].width = w
        for ci, lbl in enumerate(["KPI / Metric", "Definition", "Details / Sub-Definition"]):
            cell = tbl.cell(0, ci)
            cell.text = lbl
            cell.fill.solid(); cell.fill.fore_color.rgb = self.DARK_BLUE
            for p in cell.text_frame.paragraphs:
                p.font.bold = True; p.font.size = Pt(9)
                p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        ALT_BG = RGBColor(0xE9, 0xF0, 0xF8)
        for ri, defn in enumerate(defs, start=1):
            bg = ALT_BG if ri % 2 == 0 else RGBColor(0xFF, 0xFF, 0xFF)
            kpi_name = defn.get("kpi_name", "")
            clean_kpi = kpi_name.split("(")[0].strip() if "(" in kpi_name else kpi_name
            for ci, txt in enumerate([clean_kpi,
                                       defn.get("definition", ""),
                                       defn.get("sub_definition", "")]):
                cell = tbl.cell(ri, ci)
                cell.text = str(txt)
                cell.fill.solid(); cell.fill.fore_color.rgb = bg
                cell.text_frame.word_wrap = True
                for p in cell.text_frame.paragraphs:
                    p.font.size = Pt(8)
                    p.font.color.rgb = (RGBColor(0x1F, 0x38, 0x64)
                                        if ci == 0 else RGBColor(0x33, 0x33, 0x33))
                    if ci == 0:
                        p.font.bold = True

    # ── RAG legend helper ─────────────────────────────────────────────────────

    def _add_rag_legend(self, slide, rag_thresholds: Optional[Dict],
                        x: float, y: float) -> None:
        """
        Adds a RAG legend in a crisp white box with coloured indicator rows.
        x, y in inches. Thresholds: {"green": 95.0, "amber": 90.0}.
        """
        if not rag_thresholds:
            return
        green = float(rag_thresholds.get("green", 95.0))
        amber = float(rag_thresholds.get("amber", 90.0))

        BOX_W, BOX_H = 2.60, 0.52
        # White background box with a thin dark border — stands out against the red header
        box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x), Inches(y), Inches(BOX_W), Inches(BOX_H))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        box.line.color.rgb = RGBColor(0x33, 0x33, 0x33)
        box.line.width = Pt(0.75)

        rows = [
            (RGBColor(0x1A, 0x73, 0x48), f"■  Green  ≥{green:.0f}%"),
            (RGBColor(0xFF, 0xA0, 0x00), f"■  Amber  {amber:.0f} – {green - 1:.0f}%"),
            (RGBColor(0xCC, 0x00, 0x00), f"■  Red    <{amber:.0f}%"),
        ]
        tb = slide.shapes.add_textbox(
            Inches(x + 0.12), Inches(y + 0.05),
            Inches(BOX_W - 0.18), Inches(BOX_H - 0.08))
        tf = tb.text_frame
        tf.word_wrap = False
        first = True
        for color, label in rows:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.space_before = Pt(0); p.space_after = Pt(0)
            r = p.add_run()
            r.text = label
            r.font.size = Pt(7.0); r.font.bold = True
            r.font.color.rgb = color

    # ── Main slide renderer ───────────────────────────────────────────────────

    def _render_group_slide(
        self,
        slide,
        sheet_name: str,
        ai_output: Dict[str, Any],
        kpi_rows: List[Dict[str, Any]],
        chart_bytes: Optional[bytes],
        kpi_definitions: Optional[List[Dict]] = None,
        rag_thresholds: Optional[Dict] = None,
        parsed_sections: Optional[Dict] = None,
        full_width: bool = False,
    ):
        """
        Left column (~60 %):
          1. KEY ACHIEVEMENTS  — deep-green header + card bullets  (Excel Achievements)
          2. NEXT WEEK FOCUS   — same green format                 (Excel Focus)
          3. CONCERNS          — red header + light body           (Excel Concerns)
        Right column (~40 %): native python-pptx chart; data-table fallback.
        Content source: parsed_sections (Excel) with AI output fallback.
        full_width=True: text column spans the full slide (no chart zone reserved).
        """
        NAVY = RGBColor(0x1F, 0x38, 0x64)

        # Apply display name mapping so visual headers can differ from sheet names.
        display_name = self._DISPLAY_NAME_MAP.get(sheet_name.lower().strip(), sheet_name)

        # ── HEADER BAR ────────────────────────────────────────────────────────
        top_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), self.prs.slide_width, Inches(self._HDR_H))
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = self.DARK_BLUE
        top_bar.line.fill.background()

        tb = slide.shapes.add_textbox(
            Inches(0.2), Inches(0.08), Inches(8.2), Inches(self._HDR_H - 0.08))
        p = tb.text_frame.paragraphs[0]
        p.text = f"{display_name}  |  {self._fmt_date(self.report_date)}"
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.font.size = Pt(14); p.font.bold = True

        # ── RAG LEGEND (team-specific thresholds shown in header) ─────────────
        self._add_rag_legend(slide, rag_thresholds, x=8.5, y=0.09)

        # ── RAG BADGE ─────────────────────────────────────────────────────────
        overall_rag = ai_output.get("overall_rag", "AMBER")
        rag_bg = {"GREEN": RGBColor(0x1A, 0x73, 0x48), "AMBER": RGBColor(0xFF, 0xC0, 0x00),
                  "RED": self.MUFG_RED}.get(overall_rag, RGBColor(0xFF, 0xC0, 0x00))
        rag_fg = {"GREEN": RGBColor(0xFF, 0xFF, 0xFF), "AMBER": RGBColor(0x33, 0x33, 0x33),
                  "RED": RGBColor(0xFF, 0xFF, 0xFF)}.get(overall_rag, RGBColor(0x33, 0x33, 0x33))
        rb = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(11.3), Inches(0.09), Inches(1.9), Inches(0.43))
        rb.fill.solid(); rb.fill.fore_color.rgb = rag_bg; rb.line.fill.background()
        rt = slide.shapes.add_textbox(Inches(11.3), Inches(0.09), Inches(1.9), Inches(0.43))
        rp = rt.text_frame.paragraphs[0]
        rp.text = f"Overall: {overall_rag}"
        rp.font.bold = True; rp.font.size = Pt(11)
        rp.font.color.rgb = rag_fg; rp.alignment = PP_ALIGN.CENTER

        # ── SECTION CONTENT — Excel-parsed; fall back to AI output ───────────
        ps = parsed_sections or {}
        ach_items    = (ps.get("achievements") or
                        self._sanitize_text_list(ai_output.get("key_achievements") or []))
        focus_items  = (ps.get("focus") or
                        self._sanitize_text_list(ai_output.get("insights") or []))
        concern_items = (ps.get("concerns") or
                         ([ai_output.get("summary")] if ai_output.get("summary") else []))

        # ── Fallback for empty concerns ───────────────────────────────────────
        if not concern_items or all(not str(c).strip() for c in concern_items):
            concern_items = [
                "No critical concerns or blockers identified for this reporting period."
            ]

        # ── LEFT COLUMN: Deterministic Grid Math ─────────────────────────────
        #
        # Hard boundaries derived purely from class constants — no shape
        # property is ever read back after creation.
        #
        # top_anchor    = perf_table.bottom + margin
        # bottom_anchor = kpi_def_table.top  - margin
        # section_block = total / 3   (equal thirds, no flex)
        #
        # Each block:   blue header at block_top
        #               white body  at block_top + HEADER_H
        #
        # Z-order: all white bodies drawn first (PASS 1), then all blue
        # headers drawn last (PASS 2) — so headers are always on top.

        left_x = Inches(0)
        # When there is no chart zone, expand the text column to fill the slide.
        # Y-coordinates are NEVER changed — only the width varies.
        slide_w_in = self.prs.slide_width.inches
        left_w = Inches(slide_w_in - 0.20) if full_width else Inches(self._LEFT_W)

        # ── Hard boundary constants (all in inches, pure arithmetic) ──────────
        # top_anchor is based on _CONTENT_Y + buffer, not _PERF_Y + _PERF_H.
        # The perf table's column headers wrap to 2–3 lines, making the table
        # render taller than the nominal _PERF_H = 0.80". _CONTENT_Y = 1.42"
        # is the spec bottom-edge; adding 0.25" clears any visual overflow.
        BOT_MARGIN    = 0.10
        HEADER_H      = 0.30   # blue strip — fixed, never measured from a shape
        BASE_PT       = 11.0

        # _CONTENT_Y = 1.42" is the nominal table bottom, but the header row wraps
        # to 2–3 lines even at Pt(9), pushing the visual bottom to ~1.7–1.8".
        # 0.50" buffer brings top_anchor to 1.92" — safely below any overflow.
        top_anchor    = self._CONTENT_Y + 0.50   # 1.92"
        bottom_anchor = self._KPI_DEF_Y - BOT_MARGIN               # 6.40"
        total_avail   = bottom_anchor - top_anchor                  # 4.48"
        block_h       = total_avail / 3                             # 1.493" each
        white_h       = block_h - HEADER_H                         # body height

        SEC_HDR_BG = RGBColor(0x1F, 0x38, 0x64)
        SEC_BG     = RGBColor(0xE4, 0xED, 0xF9)
        SEC_TXT    = RGBColor(0x1F, 0x38, 0x64)
        SEC_NUM    = RGBColor(0x0D, 0x47, 0xA1)
        BORDER_CLR = RGBColor(0x1F, 0x38, 0x64)

        # Section definitions: (slot_index, items, label, bullet_char)
        section_defs = [
            (0, ach_items,     "KEY ACHIEVEMENTS", "✓"),
            (1, focus_items,   "NEXT WEEK FOCUS",  "✓"),
            (2, concern_items, "CONCERNS",         "▸"),
        ]

        # ── PASS 1: white background rects + body textboxes (low z-order) ────
        # All body shapes go into the XML before any header, so headers will
        # always paint on top regardless of coordinate proximity.
        for idx, items, label, bullet_char in section_defs:
            block_top  = top_anchor + idx * block_h   # exact grid slot
            body_top   = block_top + HEADER_H          # strictly below header

            # Full-block background rectangle (light blue, bordered)
            bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                left_x, Inches(block_top), left_w, Inches(block_h))
            bg.fill.solid()
            bg.fill.fore_color.rgb = SEC_BG
            bg.line.color.rgb = BORDER_CLR
            bg.line.width = Pt(1)

            # White body textbox — top is block_top + HEADER_H (pure math)
            body_tb = slide.shapes.add_textbox(
                left_x + Inches(0.15),
                Inches(body_top),
                left_w - Inches(0.25),
                Inches(max(white_h - 0.05, 0.20)))
            body_tb.line.color.rgb = BORDER_CLR
            body_tb.line.width = Pt(1)
            btf = body_tb.text_frame
            btf.word_wrap = True
            btf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            first = True
            for item in (items or [])[:4]:
                p = btf.paragraphs[0] if first else btf.add_paragraph()
                first = False
                self._apply_rich_bullet_to_para(
                    p, str(item), bullet_char=bullet_char,
                    base_pt=BASE_PT, text_color=SEC_TXT, num_color=SEC_NUM)

        # ── PASS 2: blue header rects + labels (high z-order) ────────────────
        # Added after all body shapes → always rendered on top in PowerPoint.
        for idx, items, label, bullet_char in section_defs:
            block_top = top_anchor + idx * block_h

            hdr_rect = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                left_x, Inches(block_top), left_w, Inches(HEADER_H))
            hdr_rect.fill.solid()
            hdr_rect.fill.fore_color.rgb = SEC_HDR_BG
            hdr_rect.line.fill.background()

            hdr_tb = slide.shapes.add_textbox(
                left_x + Inches(0.14),
                Inches(block_top + 0.05),
                left_w - Inches(0.20),
                Inches(HEADER_H - 0.05))
            hp = hdr_tb.text_frame.paragraphs[0]
            hp.text = label
            hp.font.bold  = True
            hp.font.size  = Pt(11)
            hp.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

        # ── RIGHT COLUMN: chart_bytes fallback (simple path only) ─────────────
        if chart_bytes:
            try:
                slide.shapes.add_picture(
                    io.BytesIO(chart_bytes),
                    Inches(self._RIGHT_X), Inches(self._CONTENT_Y),
                    width=Inches(self._CHART_W))
            except Exception:
                pass

        self._add_kpi_definition_table(slide, kpi_definitions or [])

    # ── Multigroup slide renderer (e.g. CMDB team: CMDB / HAM / SAM) ─────────

    def _render_multigroup_slide(
        self,
        slide,
        sheet_name: str,
        ai_output: Dict[str, Any],
        kpi_rows: List[Dict[str, Any]],
        kpi_definitions: Optional[List[Dict]] = None,
        multigroup_data: Optional[Dict] = None,
        rag_thresholds: Optional[Dict] = None,
        df: Optional[pd.DataFrame] = None,
    ) -> None:
        """
        Special full-width slide for teams with 3 sub-groups (no charts).
        Layout: header bar → KPI table → 3-column body → timeline → KPI def table.
        """
        NAVY = RGBColor(0x1F, 0x38, 0x64)
        COL_HEADER_COLORS = [
            RGBColor(0x1F, 0x38, 0x64),  # dark navy  (CMDB)
            RGBColor(0x0F, 0x6E, 0x56),  # teal       (HAM)
            RGBColor(0x8B, 0x45, 0x00),  # warm amber (3rd group)
        ]
        SEC_HDR_COLORS = [
            RGBColor(0x2E, 0x4F, 0x8C),
            RGBColor(0x17, 0x8A, 0x6C),
            RGBColor(0xAA, 0x5A, 0x14),
        ]
        SEC_BG_COLORS = [
            [RGBColor(0xED, 0xFB, 0xF3), RGBColor(0xFF, 0xF0, 0xF0), RGBColor(0xF0, 0xF5, 0xFF)],
            [RGBColor(0xED, 0xFB, 0xF3), RGBColor(0xFF, 0xF0, 0xF0), RGBColor(0xF0, 0xF5, 0xFF)],
            [RGBColor(0xFD, 0xF9, 0xEC), RGBColor(0xFF, 0xF0, 0xF0), RGBColor(0xF5, 0xF0, 0xFF)],
        ]

        # ── HEADER BAR ─────────────────────────────────────────────────────
        top_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), self.prs.slide_width, Inches(self._HDR_H))
        top_bar.fill.solid(); top_bar.fill.fore_color.rgb = self.DARK_BLUE
        top_bar.line.fill.background()

        tb = slide.shapes.add_textbox(
            Inches(0.2), Inches(0.08), Inches(8.2), Inches(self._HDR_H - 0.08))
        p = tb.text_frame.paragraphs[0]
        p.text = f"{sheet_name}  |  {self._fmt_date(self.report_date)}"
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.font.size = Pt(14); p.font.bold = True

        self._add_rag_legend(slide, rag_thresholds, x=8.5, y=0.09)

        # ── RAG BADGE ──────────────────────────────────────────────────────
        overall_rag = ai_output.get("overall_rag", "AMBER")
        rag_bg = {"GREEN": RGBColor(0x1A, 0x73, 0x48), "AMBER": RGBColor(0xFF, 0xC0, 0x00),
                  "RED": self.MUFG_RED}.get(overall_rag, RGBColor(0xFF, 0xC0, 0x00))
        rag_fg = {"GREEN": RGBColor(0xFF, 0xFF, 0xFF), "AMBER": RGBColor(0x33, 0x33, 0x33),
                  "RED": RGBColor(0xFF, 0xFF, 0xFF)}.get(overall_rag, RGBColor(0x33, 0x33, 0x33))
        rb = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(11.3), Inches(0.09), Inches(1.9), Inches(0.43))
        rb.fill.solid(); rb.fill.fore_color.rgb = rag_bg; rb.line.fill.background()
        rt = slide.shapes.add_textbox(Inches(11.3), Inches(0.09), Inches(1.9), Inches(0.43))
        rp = rt.text_frame.paragraphs[0]
        rp.text = f"Overall: {overall_rag}"
        rp.font.bold = True; rp.font.size = Pt(11)
        rp.font.color.rgb = rag_fg; rp.alignment = PP_ALIGN.CENTER

        # ── KPI PERFORMANCE TABLE (top strip, full-width) ───────────────────
        if df is not None and not df.empty:
            MG_PREFIXES = ("Group1_", "Group2_", "Group3_", "Timeline")
            kpi_df = df[[c for c in df.columns
                         if not any(str(c).startswith(p) for p in MG_PREFIXES)]]
            self._add_performance_table(slide, kpi_df, kpi_rows)

        # ── LAYOUT CONSTANTS ───────────────────────────────────────────────
        BODY_Y      = self._CONTENT_Y          # 1.07"
        GRP_HDR_H   = 0.33
        TIMELINE_H  = 0.68
        TIMELINE_GAP = 0.08
        KPI_DEF_Y   = self._KPI_DEF_Y          # 6.50"
        BODY_H      = KPI_DEF_Y - BODY_Y - TIMELINE_H - TIMELINE_GAP  # ≈4.62"
        SEC_COUNT   = 3
        SEC_H       = (BODY_H - GRP_HDR_H) / SEC_COUNT               # ≈1.43"
        SEC_HDR_H   = 0.27
        N_SLIDE_W   = 13.33

        groups = (multigroup_data or {}).get("groups", [])
        n_cols  = max(1, min(3, len(groups)))
        col_w   = N_SLIDE_W / n_cols

        # ── 3-COLUMN BODY ───────────────────────────────────────────────────
        for gi, group in enumerate(groups[:3]):
            gx        = gi * col_w
            col_color = COL_HEADER_COLORS[gi % len(COL_HEADER_COLORS)]
            sec_hdr_c = SEC_HDR_COLORS[gi % len(SEC_HDR_COLORS)]
            sec_bgs   = SEC_BG_COLORS[gi % len(SEC_BG_COLORS)]
            group_name = group.get("name", f"Group {gi + 1}")

            # Group column header
            gh = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(gx), Inches(BODY_Y), Inches(col_w), Inches(GRP_HDR_H))
            gh.fill.solid(); gh.fill.fore_color.rgb = col_color; gh.line.fill.background()
            gt = slide.shapes.add_textbox(
                Inches(gx + 0.12), Inches(BODY_Y + 0.07),
                Inches(col_w - 0.15), Inches(GRP_HDR_H - 0.10))
            gp = gt.text_frame.paragraphs[0]
            gp.text = group_name.upper()
            gp.font.bold = True; gp.font.size = Pt(12)
            gp.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

            sections = [
                ("Achievements",    group.get("achievements", "")),
                ("Concerns",        group.get("concerns", "")),
                ("Next Week Focus",  group.get("focus", "")),
            ]
            for si, (sec_label, sec_text) in enumerate(sections):
                sy   = BODY_Y + GRP_HDR_H + si * SEC_H
                sbg  = sec_bgs[si]

                # Section background
                sb = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(gx), Inches(sy), Inches(col_w), Inches(SEC_H))
                sb.fill.solid(); sb.fill.fore_color.rgb = sbg
                sb.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC); sb.line.width = Pt(0.5)

                # Section label bar
                shb = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(gx), Inches(sy), Inches(col_w), Inches(SEC_HDR_H))
                shb.fill.solid(); shb.fill.fore_color.rgb = sec_hdr_c
                shb.line.fill.background()
                sht = slide.shapes.add_textbox(
                    Inches(gx + 0.12), Inches(sy + 0.05),
                    Inches(col_w - 0.15), Inches(SEC_HDR_H - 0.07))
                shp = sht.text_frame.paragraphs[0]
                shp.text = sec_label.upper()
                shp.font.bold = True; shp.font.size = Pt(8.5)
                shp.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

                # Section content (bullet points)
                if sec_text and str(sec_text).strip():
                    stb = slide.shapes.add_textbox(
                        Inches(gx + 0.12), Inches(sy + SEC_HDR_H + 0.07),
                        Inches(col_w - 0.22), Inches(SEC_H - SEC_HDR_H - 0.10))
                    stf = stb.text_frame
                    stf.word_wrap = True
                    stf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
                    bullets = [b.strip() for b in re.split(r"[;\n]", str(sec_text)) if b.strip()]
                    for bi, bullet in enumerate(bullets[:4]):
                        bp = stf.paragraphs[0] if bi == 0 else stf.add_paragraph()
                        self._apply_rich_bullet_to_para(
                            bp, bullet, bullet_char="▸", base_pt=7.5,
                            text_color=RGBColor(0x22, 0x22, 0x22),
                            num_color=col_color)

        # ── TIMELINE STRIP (full-width, below 3-column body) ────────────────
        timeline_y    = BODY_Y + BODY_H + TIMELINE_GAP
        timeline_text = (multigroup_data or {}).get("timeline", "")

        tlbg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(timeline_y), self.prs.slide_width, Inches(TIMELINE_H))
        tlbg.fill.solid(); tlbg.fill.fore_color.rgb = RGBColor(0xF2, 0xF5, 0xFA)
        tlbg.line.color.rgb = NAVY; tlbg.line.width = Pt(0.75)

        tlhb = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(timeline_y), self.prs.slide_width, Inches(0.27))
        tlhb.fill.solid(); tlhb.fill.fore_color.rgb = NAVY; tlhb.line.fill.background()
        tlht = slide.shapes.add_textbox(
            Inches(0.14), Inches(timeline_y + 0.04), Inches(6), Inches(0.22))
        tlhp = tlht.text_frame.paragraphs[0]
        tlhp.text = "TIMELINE  &  KEY UPDATES"
        tlhp.font.bold = True; tlhp.font.size = Pt(9)
        tlhp.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

        if timeline_text and str(timeline_text).strip():
            tltb = slide.shapes.add_textbox(
                Inches(0.14), Inches(timeline_y + 0.30),
                self.prs.slide_width - Inches(0.28), Inches(TIMELINE_H - 0.34))
            tltf = tltb.text_frame
            tltf.word_wrap = True
            tltf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
            tltp = tltf.paragraphs[0]
            tltp.text = str(timeline_text)
            tltp.font.size = Pt(8.5)
            tltp.font.color.rgb = RGBColor(0x22, 0x22, 0x22)

        # ── KPI DEFINITION TABLE ────────────────────────────────────────────
        self._add_kpi_definition_table(slide, kpi_definitions or [])

    # ── Native chart helpers ──────────────────────────────────────────────────

    def _add_line_overlay_to_chart(self, chart, cat_ax_id: str,
                                    primary_val_ax_id: str) -> None:
        """
        Move the last bar series into a new lineChart with its own secondary
        Y-axis (right side).  This prevents scale conflicts when bars show counts
        and the line shows percentages (or vice versa).
        """
        try:
            from lxml import etree
            C = "http://schemas.openxmlformats.org/drawingml/2006/chart"
            A = "http://schemas.openxmlformats.org/drawingml/2006/main"

            plot_area = chart._element.find(f".//{{{C}}}plotArea")
            bar_chart = plot_area.find(f"{{{C}}}barChart") if plot_area is not None else None
            if bar_chart is None:
                return

            all_sers = bar_chart.findall(f"{{{C}}}ser")
            if len(all_sers) < 2:
                return

            line_ser = all_sers[-1]
            bar_chart.remove(line_ser)

            # Secondary value axis ID = primary + 1 (guaranteed unique)
            sec_ax_id = str(int(primary_val_ax_id) + 1)

            # ── Style the line series ──────────────────────────────────────
            spPr = line_ser.find(f"{{{C}}}spPr")
            if spPr is None:
                spPr = etree.SubElement(line_ser, f"{{{C}}}spPr")
            for child in list(spPr):
                spPr.remove(child)
            ln = etree.SubElement(spPr, f"{{{A}}}ln")
            ln.set("w", "22860")   # 1.8 pt
            sf = etree.SubElement(ln, f"{{{A}}}solidFill")
            etree.SubElement(sf, f"{{{A}}}srgbClr").set("val", "1F3864")

            marker = etree.SubElement(line_ser, f"{{{C}}}marker")
            etree.SubElement(marker, f"{{{C}}}symbol").set("val", "circle")
            etree.SubElement(marker, f"{{{C}}}size").set("val", "5")
            etree.SubElement(line_ser, f"{{{C}}}smooth").set("val", "0")

            # ── Build lineChart element ────────────────────────────────────
            lc = etree.Element(f"{{{C}}}lineChart")
            etree.SubElement(lc, f"{{{C}}}grouping").set("val", "standard")
            etree.SubElement(lc, f"{{{C}}}varyColors").set("val", "0")
            lc.append(line_ser)
            # Share cat axis; use new secondary val axis
            etree.SubElement(lc, f"{{{C}}}axId").set("val", cat_ax_id)
            etree.SubElement(lc, f"{{{C}}}axId").set("val", sec_ax_id)

            bar_idx = list(plot_area).index(bar_chart)
            plot_area.insert(bar_idx + 1, lc)

            # ── Add secondary valAx element ───────────────────────────────
            sec_val_ax = etree.SubElement(plot_area, f"{{{C}}}valAx")
            etree.SubElement(sec_val_ax, f"{{{C}}}axId").set("val", sec_ax_id)

            scaling = etree.SubElement(sec_val_ax, f"{{{C}}}scaling")
            etree.SubElement(scaling, f"{{{C}}}orientation").set("val", "minMax")

            etree.SubElement(sec_val_ax, f"{{{C}}}delete").set("val", "0")
            etree.SubElement(sec_val_ax, f"{{{C}}}axPos").set("val", "r")

            numFmt = etree.SubElement(sec_val_ax, f"{{{C}}}numFmt")
            numFmt.set("formatCode", "General")
            numFmt.set("sourceLinked", "0")

            etree.SubElement(sec_val_ax, f"{{{C}}}majorTickMark").set("val", "none")
            etree.SubElement(sec_val_ax, f"{{{C}}}minorTickMark").set("val", "none")
            etree.SubElement(sec_val_ax, f"{{{C}}}tickLblPos").set("val", "nextTo")
            etree.SubElement(sec_val_ax, f"{{{C}}}crossAx").set("val", cat_ax_id)
            etree.SubElement(sec_val_ax, f"{{{C}}}crosses").set("val", "max")
            etree.SubElement(sec_val_ax, f"{{{C}}}crossBetween").set("val", "midCat")

        except Exception:
            pass

    def _enable_chart_data_labels(self, chart) -> None:
        """
        Inject show-value data labels on every series.
        Bars get 'outEnd' (above bar); lines get 't' (above point).
        Font capped at 7pt to avoid overlap.
        """
        try:
            from lxml import etree
            C = "http://schemas.openxmlformats.org/drawingml/2006/chart"
            A = "http://schemas.openxmlformats.org/drawingml/2006/main"

            plot_area = chart._element.find(f".//{{{C}}}plotArea")
            if plot_area is None:
                return

            for chart_el in list(plot_area):
                local = chart_el.tag.split("}")[-1] if "}" in chart_el.tag else chart_el.tag
                is_line = (local == "lineChart")
                # bestFit lets PowerPoint auto-place line labels avoiding overflow;
                # inEnd keeps bar labels inside the bar top clear of the line above.
                pos = "bestFit" if is_line else "inEnd"

                for ser in chart_el.findall(f"{{{C}}}ser"):
                    if ser.find(f"{{{C}}}dLbls") is not None:
                        continue

                    # Build dLbls as a standalone element first
                    dLbls = etree.Element(f"{{{C}}}dLbls")

                    numFmt = etree.SubElement(dLbls, f"{{{C}}}numFmt")
                    numFmt.set("formatCode", "General")
                    numFmt.set("sourceLinked", "0")

                    txPr = etree.SubElement(dLbls, f"{{{C}}}txPr")
                    etree.SubElement(txPr, f"{{{A}}}bodyPr")
                    etree.SubElement(txPr, f"{{{A}}}lstStyle")
                    _p = etree.SubElement(txPr, f"{{{A}}}p")
                    _pPr = etree.SubElement(_p, f"{{{A}}}pPr")
                    defRPr = etree.SubElement(_pPr, f"{{{A}}}defRPr")
                    defRPr.set("sz", "550")

                    etree.SubElement(dLbls, f"{{{C}}}dLblPos").set("val", pos)
                    for tag, val in [
                        ("showLegendKey",  "0"),
                        ("showVal",        "1"),
                        ("showCatName",    "0"),
                        ("showSerName",    "0"),
                        ("showPercent",    "0"),
                        ("showBubbleSize", "0"),
                    ]:
                        etree.SubElement(dLbls, f"{{{C}}}{tag}").set("val", val)

                    # OOXML schema requires dLbls to appear BEFORE cat/val.
                    # etree.SubElement appends at the end (after cat/val) → invalid XML → corrupt PPTX.
                    # Use insert at the position of the cat element instead.
                    cat_el = ser.find(f"{{{C}}}cat")
                    if cat_el is not None:
                        ser.insert(list(ser).index(cat_el), dLbls)
                    else:
                        ser.append(dLbls)
        except Exception:
            pass

    def _hide_y_axis_labels(self, chart) -> None:
        """
        Hide tick-label numbers on both primary and secondary value axes.
        Data labels on each series make axis scale numbers redundant clutter.
        """
        try:
            from lxml import etree
            C = "http://schemas.openxmlformats.org/drawingml/2006/chart"
            plot_area = chart._element.find(f".//{{{C}}}plotArea")
            if plot_area is None:
                return
            for val_ax in plot_area.findall(f"{{{C}}}valAx"):
                for tag in ("tickLblPos", "majorTickMark", "minorTickMark"):
                    el = val_ax.find(f"{{{C}}}{tag}")
                    if el is not None:   # Only modify existing — creating new elements
                        el.set("val", "none")  # at wrong position corrupts the file
        except Exception:
            pass

    def _set_legend_style(self, chart) -> None:
        try:
            from pptx.enum.chart import XL_LEGEND_POSITION
            chart.has_legend = True
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            chart.legend.include_in_layout = False
        except Exception:
            pass

    def _try_native_chart(
        self,
        slide,
        df: pd.DataFrame,
        chart_type: str,
        column_types: dict,
        left_x: float,
        top_y: float,
        width_in: float,
        height_in: float,
        cols_override: Optional[List] = None,
    ) -> bool:
        """
        Build a 100 % native python-pptx chart using only the python-pptx API.
        No lxml XML manipulation — that caused PowerPoint repair/corruption errors.

        chart_type == "line"  → LINE_MARKERS (rate/% series)
        chart_type == "auto"  → COLUMN_CLUSTERED (count/volume series)
        cols_override         → restrict series to this column subset

        Returns True on success; False → caller shows data table instead.
        """
        if chart_type == "none" or df is None or df.empty:
            return False
        try:
            from pptx.chart.data import ChartData
            from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

            skip_kw = {"comment", "achievement", "weekly", "month", "quarter",
                       "sprint no", "week no", "date"}
            num_types = {"numeric", "percent", "percent_decimal"}

            if cols_override is not None:
                override_set = {str(c) for c in cols_override}
                numeric_cols = [c for c in df.columns
                                if str(c) in override_set
                                and column_types.get(str(c)) in num_types]
            else:
                numeric_cols = [
                    c for c in df.columns
                    if column_types.get(str(c)) in num_types
                    and not any(kw in str(c).lower() for kw in skip_kw)
                ]
            if not numeric_cols:
                return False

            cat_col = next(
                (c for c in df.columns
                 if column_types.get(str(c)) == "text"
                 and not any(kw in str(c).lower() for kw in skip_kw)),
                None,
            )
            cats = ([str(v) if v and str(v).lower() != "nan" else f"R{i+1}"
                     for i, v in enumerate(df[cat_col])]
                    if cat_col is not None
                    else [f"W{i+1}" for i in range(len(df))])

            chart_data = ChartData()
            chart_data.categories = cats

            # Cap at 3 series for readability; truncate long legend names to 18 chars
            for col in numeric_cols[:3]:
                vals = []
                for v in df[col]:
                    try:
                        fv = float(v)
                        if column_types.get(str(col)) == "percent_decimal":
                            fv *= 100
                        vals.append(round(fv, 2))
                    except Exception:
                        vals.append(None)
                raw_name = str(col).split("(")[0].strip()
                short_name = raw_name if len(raw_name) <= 18 else raw_name[:16] + "…"
                chart_data.add_series(short_name, vals)

            n_series = min(len(numeric_cols), 3)
            use_line = (chart_type == "line")
            xl_type  = XL_CHART_TYPE.LINE_MARKERS if use_line else XL_CHART_TYPE.COLUMN_CLUSTERED

            chart_frame = slide.shapes.add_chart(
                xl_type,
                Inches(left_x), Inches(top_y),
                Inches(width_in), Inches(height_in),
                chart_data,
            )
            chart = chart_frame.chart
            chart.has_title = False

            MUFG = [
                RGBColor(0xCC, 0x00, 0x00),
                RGBColor(0x1F, 0x38, 0x64),
                RGBColor(0x59, 0x57, 0x57),
            ]
            for si, series in enumerate(chart.series):
                series.format.fill.solid()
                series.format.fill.fore_color.rgb = MUFG[si % len(MUFG)]
                try:
                    series.data_labels.show_value = True
                    if use_line:
                        # Line charts: alternate ABOVE/BELOW so labels never overlap
                        series.data_labels.font.size = Pt(8)
                        try:
                            from pptx.enum.chart import XL_DATA_LABEL_POSITION
                            series.data_labels.position = (
                                XL_DATA_LABEL_POSITION.ABOVE
                                if si % 2 == 0
                                else XL_DATA_LABEL_POSITION.BELOW
                            )
                        except Exception:
                            pass
                    else:
                        series.data_labels.font.size = Pt(7)
                except Exception:
                    pass

            # Remove Y-axis labels and all gridlines — data labels are sufficient
            try:
                chart.value_axis.visible = False
            except Exception:
                pass
            try:
                chart.value_axis.has_major_gridlines = False
                chart.value_axis.has_minor_gridlines = False
            except Exception:
                pass
            try:
                chart.category_axis.has_major_gridlines = False
            except Exception:
                pass

            # Legend at bottom with compact 9pt font so bars dominate
            if n_series > 1:
                chart.has_legend = True
                chart.legend.position = XL_LEGEND_POSITION.BOTTOM
                chart.legend.include_in_layout = False
                try:
                    chart.legend.font.size = Pt(9)
                except Exception:
                    pass
            else:
                chart.has_legend = False

            # Maximize plot area so bars/lines fill the chart frame
            try:
                from pptx.util import Emu
                pa = chart.plot_area
                pa.left   = Emu(0)
                pa.top    = Emu(0)
                pa.width  = None   # let pptx auto-size to frame width
                pa.height = None
            except Exception:
                pass

            # Bounding-box border: transparent rectangle drawn on top of the chart
            # with the identical coordinates — python-pptx line borders on charts
            # are unreliable; this AutoShape approach always works.
            try:
                bf = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(left_x), Inches(top_y),
                    Inches(width_in), Inches(height_in),
                )
                bf.fill.background()
                bf.line.fill.solid()
                bf.line.color.rgb = RGBColor(0xA6, 0xA6, 0xA6)
                bf.line.width = Pt(1)
            except Exception:
                pass

            return True
        except Exception:
            return False

    def _try_native_pie_chart(
        self,
        slide,
        df: pd.DataFrame,
        column_types: dict,
        left_x: float,
        top_y: float,
        width_in: float,
        height_in: float,
        cols_override: Optional[List] = None,
    ) -> bool:
        """
        Render a PIE chart from the LAST row of count (numeric) columns.
        Uses XL_CHART_TYPE.PIE with percentage data labels.
        Suitable when there are 2-5 count columns and a single time period matters.
        Returns True on success; False → caller falls back.
        """
        if df is None or df.empty:
            return False
        try:
            from pptx.chart.data import CategoryChartData
            from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

            skip_kw = {"comment", "achievement", "weekly", "month", "quarter",
                       "sprint no", "week no", "date"}
            num_types = {"numeric", "percent", "percent_decimal"}

            if cols_override is not None:
                pie_cols = [c for c in cols_override
                            if column_types.get(str(c)) in num_types]
            else:
                pie_cols = [c for c in df.columns
                            if column_types.get(str(c)) in num_types
                            and not any(kw in str(c).lower() for kw in skip_kw)]

            if len(pie_cols) < 2:
                return False

            last_row = df.iloc[-1]
            # Cap at top 5 slices by absolute value when more columns are available
            if len(pie_cols) > 5:
                def _safe_val(c):
                    try:
                        return abs(float(last_row[c]))
                    except Exception:
                        return 0.0
                pie_cols = sorted(pie_cols, key=_safe_val, reverse=True)[:5]

            slices, labels = [], []
            for col in pie_cols:
                try:
                    raw_val = float(last_row[col])
                    # Normalise percent_decimal (0.97 → 97) so slices are comparable
                    if column_types.get(str(col)) == "percent_decimal":
                        raw_val = raw_val * 100
                    if raw_val > 0:
                        raw = str(col).split("(")[0].strip()
                        labels.append(raw if len(raw) <= 20 else raw[:18] + "…")
                        slices.append(round(raw_val, 2))
                except Exception:
                    pass

            if len(slices) < 2:
                return False

            chart_data = CategoryChartData()
            chart_data.categories = labels
            chart_data.add_series("Values", tuple(slices))

            chart_frame = slide.shapes.add_chart(
                XL_CHART_TYPE.PIE,
                Inches(left_x), Inches(top_y),
                Inches(width_in), Inches(height_in),
                chart_data,
            )
            chart = chart_frame.chart
            chart.has_title = False
            chart.has_legend = True
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            chart.legend.include_in_layout = False
            try:
                chart.legend.font.size = Pt(9)
            except Exception:
                pass

            # Series-level data labels — more reliable than plot-level for pie
            try:
                series = chart.plots[0].series[0]
                series.has_data_labels = True
                series.data_labels.show_percentage = True
                series.data_labels.show_category_name = True
                series.data_labels.show_value = False
                try:
                    series.data_labels.font.size = Pt(10)
                except Exception:
                    pass
            except Exception:
                # Fallback: plot-level labels
                try:
                    dls = chart.plots[0].data_labels
                    dls.show_percentage = True
                    dls.show_value = False
                    dls.font.size = Pt(9)
                except Exception:
                    pass

            # Bounding-box border — same coords as chart, transparent fill
            try:
                bf = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(left_x), Inches(top_y),
                    Inches(width_in), Inches(height_in),
                )
                bf.fill.background()
                bf.line.fill.solid()
                bf.line.color.rgb = RGBColor(0xA6, 0xA6, 0xA6)
                bf.line.width = Pt(1)
            except Exception:
                pass

            return True
        except Exception:
            return False

    # ── Slide management helpers ──────────────────────────────────────────────

    def _remove_slide(self, slide) -> None:
        sldIdLst = self.prs.slides._sldIdLst
        for s, sl in zip(self.prs.slides, list(sldIdLst)):
            if s is slide:
                self.prs.part.drop_rel(sl.rId)
                sldIdLst.remove(sl)
                return

    def _remove_duplicate_slides(self, sheet_name: str, keep_slide) -> None:
        prefix = f"{sheet_name}  |  "
        to_remove = []
        for slide in list(self.prs.slides):
            if slide is keep_slide:
                continue
            for shape in slide.shapes:
                if not getattr(shape, "has_text_frame", False):
                    continue
                if (shape.text_frame.text or "").strip().startswith(prefix):
                    to_remove.append(slide)
                    break
        for slide in to_remove:
            self._remove_slide(slide)

    def add_title_slide(self, report_date: str, group_count: int):
        self.report_date = report_date
        slide = (self.prs.slides[0]
                 if len(self.prs.slides) > 0
                 else self.prs.slides.add_slide(
                     self.prs.slide_layouts[self._blank_slide_layout()]))
        self._clear_slide(slide)

        top_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
            self.prs.slide_width, Inches(1))
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = self.DARK_BLUE
        top_bar.line.fill.background()

        tb = slide.shapes.add_textbox(Inches(0), Inches(1.1), self.prs.slide_width, Inches(1.0))
        tf = tb.text_frame; tf.clear()
        p = tf.paragraphs[0]
        p.text = f"ITSM Weekly Report — {self._fmt_date(report_date)}"
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Inches(0.35); p.font.bold = True
        p.font.color.rgb = self.MUFG_RED

        tb2 = slide.shapes.add_textbox(Inches(0), Inches(2.0), self.prs.slide_width, Inches(0.6))
        p2 = tb2.text_frame.paragraphs[0]
        p2.text = f"Auto-generated | {group_count} groups | AI-powered insights"
        p2.alignment = PP_ALIGN.CENTER
        p2.font.size = Inches(0.18); p2.font.color.rgb = self.MUFG_GRAY

        tb3 = slide.shapes.add_textbox(Inches(0), Inches(6.9), self.prs.slide_width, Inches(0.4))
        p3 = tb3.text_frame.paragraphs[0]
        p3.text = "Generated by ITSM Automation Agent | Confidential"
        p3.alignment = PP_ALIGN.CENTER
        p3.font.size = Inches(0.13); p3.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    def add_group_slide(
        self,
        sheet_name: str,
        ai_output: Dict[str, Any],
        kpi_rows: List[Dict[str, Any]],
        chart_bytes: Optional[bytes],
        charts: Optional[List[Dict]] = None,
        summary_mode: str = "ai_write",
        df: Optional[pd.DataFrame] = None,
        kpi_definitions: Optional[List[Dict]] = None,
        column_types: Optional[Dict] = None,
        layout: str = "standard",
        multigroup_data: Optional[Dict] = None,
        rag_thresholds: Optional[Dict] = None,
        parsed_sections: Optional[Dict] = None,
    ):
        if df is None or (isinstance(df, pd.DataFrame) and df.empty):
            return

        existing = self._find_group_slide(sheet_name)
        if existing is None:
            slide = self.prs.slides.add_slide(
                self.prs.slide_layouts[self._blank_slide_layout()])
        else:
            slide = existing
            self._clear_slide(slide)
            self._remove_duplicate_slides(sheet_name, keep_slide=slide)

        if layout == "multigroup":
            self._render_multigroup_slide(
                slide, sheet_name, ai_output, kpi_rows,
                kpi_definitions or [], multigroup_data or {},
                rag_thresholds, df)
        elif charts:
            self._render_group_slide_multi(
                slide, sheet_name, ai_output, kpi_rows,
                charts, summary_mode, df, kpi_definitions or [],
                column_types or {}, rag_thresholds, parsed_sections)
        else:
            # Draw perf table FIRST so section headers have higher z-order.
            # This ensures the blue headers are never hidden by the table even
            # if the table's rendered height overflows its nominal _PERF_H bound.
            self._add_performance_table(slide, df, kpi_rows)
            # full_width=True when there is no chart zone — text spans the slide.
            no_charts = not chart_bytes
            self._render_group_slide(slide, sheet_name, ai_output, kpi_rows,
                                     chart_bytes, kpi_definitions or [],
                                     rag_thresholds, parsed_sections,
                                     full_width=no_charts)
            if summary_mode == "use_excel":
                self._add_team_provided_label(slide, Inches(0.2), Inches(2.4))


    def _add_team_provided_label(self, slide, left, top):
        tb = slide.shapes.add_textbox(left, top, Inches(3), Inches(0.25))
        p = tb.text_frame.paragraphs[0]
        p.text = "(Team-provided)"
        p.font.italic = True; p.font.size = Pt(8)
        p.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

    def _add_data_table(self, slide, df: pd.DataFrame, zone_coords: tuple) -> None:
        """Render a compact python-pptx table as chart-zone fallback. Max 4 rows, 8 cols."""
        if df is None or df.empty:
            return
        display_df = df.dropna(how="all").tail(4)
        cols = list(display_df.columns)[:8]
        display_df = display_df[cols]
        n_rows, n_cols = len(display_df) + 1, len(cols)
        if n_rows < 2 or n_cols < 1:
            return
        left_x, top_y, width, height = zone_coords
        tbl_shape = slide.shapes.add_table(
            n_rows, n_cols,
            Inches(left_x), Inches(top_y), Inches(width), Inches(height))
        tbl = tbl_shape.table
        for ci, col_name in enumerate(cols):
            cell = tbl.cell(0, ci)
            cell.text = str(col_name)
            cell.fill.solid(); cell.fill.fore_color.rgb = self.DARK_BLUE
            for p in cell.text_frame.paragraphs:
                p.font.bold = True; p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                p.font.size = Pt(8); p.alignment = PP_ALIGN.CENTER
        for ri, (_, row) in enumerate(display_df.iterrows(), start=1):
            bg = self.LIGHT_BG if ri % 2 == 0 else RGBColor(0xFF, 0xFF, 0xFF)
            for ci, col_name in enumerate(cols):
                cell = tbl.cell(ri, ci)
                val = row[col_name]
                cell.text = ("" if val is None or
                              (isinstance(val, float) and val != val) else str(val))
                cell.fill.solid(); cell.fill.fore_color.rgb = bg
                for p in cell.text_frame.paragraphs:
                    p.font.size = Pt(7)
                    try:
                        float(val); p.alignment = PP_ALIGN.RIGHT
                    except (TypeError, ValueError):
                        p.alignment = PP_ALIGN.LEFT

    def _render_group_slide_multi(
        self,
        slide,
        sheet_name: str,
        ai_output: Dict[str, Any],
        kpi_rows: List[Dict[str, Any]],
        charts: List[Dict],
        summary_mode: str,
        df: Optional[pd.DataFrame],
        kpi_definitions: Optional[List[Dict]] = None,
        column_types: Optional[Dict] = None,
        rag_thresholds: Optional[Dict] = None,
        parsed_sections: Optional[Dict] = None,
    ) -> None:
        """
        Multi-chart path.
        Left column: rendered by _render_group_slide (no chart_bytes).
        Right column: always native python-pptx chart; data table if native fails.
        No matplotlib PNG fallback — charts are 100 % editable.
        """
        self._render_group_slide(slide, sheet_name, ai_output, kpi_rows,
                                 chart_bytes=None, kpi_definitions=kpi_definitions,
                                 rag_thresholds=rag_thresholds,
                                 parsed_sections=parsed_sections,
                                 full_width=False)
        perf_tbl = self._add_performance_table(slide, df, kpi_rows)

        if summary_mode == "use_excel":
            self._add_team_provided_label(
                slide, Inches(0.2), Inches(self._CONTENT_Y + 0.9))

        # ── Dynamic chart zone: anchored to the actual perf table bottom ────────
        # Using shape.top + shape.height converts shape coordinates (EMU) back to
        # inches so charts never overlap the perf table regardless of row height.
        _EMU = 914400.0  # EMU per inch
        if perf_tbl is not None:
            chart_start_y = (perf_tbl.top + perf_tbl.height) / _EMU + 0.30
        else:
            chart_start_y = self._PERF_Y + self._PERF_H + 0.30  # fallback constant

        # Bottom anchor = top of the KPI definition table minus margin
        chart_max_bottom  = self._KPI_DEF_Y - 0.15
        available_h       = chart_max_bottom - chart_start_y

        GAP = 0.10  # vertical gap between two stacked charts

        # Config-driven chart placement — use the charts list (already capped at 2 by t5/t2).
        ct = column_types or {}
        valid_charts = [c for c in (charts or []) if c is not None][:2]

        PADDING_IN = 0.10  # inward padding so chart never touches table boundaries

        def _place_one_chart(entry: dict, top_y: float, height: float) -> None:
            """Route to the correct native chart type.
            Pie: MUST render as a pie chart — no table fallback.
            Other types: fall back to data table if native chart fails.
            Padding is applied here so both the chart and its bounding box
            sit cleanly inside the allocated slot without touching neighbours.
            """
            px = self._RIGHT_X + PADDING_IN
            py = top_y + PADDING_IN
            pw = self._CHART_W - 2 * PADDING_IN
            ph = height - 2 * PADDING_IN
            req = str(entry.get("chart_type", "auto")).strip().lower()
            if req == "pie":
                self._try_native_pie_chart(slide, df, ct, px, py, pw, ph)
            elif req in ("line", "bar_dotted_line"):
                ok = self._try_native_chart(slide, df, "line", ct, px, py, pw, ph)
                if not ok:
                    self._add_data_table(slide, df, (px, py, pw, ph))
            else:
                ok = self._try_native_chart(slide, df, "auto", ct, px, py, pw, ph)
                if not ok:
                    self._add_data_table(slide, df, (px, py, pw, ph))

        if not valid_charts:
            pass  # no charts — leave right zone blank

        elif len(valid_charts) == 1:
            # Single chart: full available height
            _place_one_chart(valid_charts[0], chart_start_y, available_h)

        else:
            # Two charts: split available height equally, with a gap between them
            chart_h = (available_h / 2) - (GAP / 2)
            y1 = chart_start_y
            y2 = chart_start_y + chart_h + GAP
            for entry, top_y in zip(valid_charts, [y1, y2]):
                _place_one_chart(entry, top_y, chart_h)

    def save(self, output_path: str) -> str:
        os.makedirs(os.path.dirname(os.path.abspath(output_path)), exist_ok=True)
        try:
            self.prs.save(output_path)
            return output_path
        except PermissionError:
            draft_path = output_path.replace(".pptx", "_draft.pptx")
            self.prs.save(draft_path)
            print(f"\n[WARNING] '{output_path}' is open in PowerPoint.\n"
                  f"          Report saved to: {draft_path}\n"
                  f"          Close PowerPoint and rename the file.")
            return draft_path
